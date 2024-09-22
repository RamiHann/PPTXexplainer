import unittest
from unittest.mock import patch, AsyncMock
from pptx import Presentation
from explainer import combine_slide_text, process_slide, process_presentations

class TestExplainer(unittest.TestCase):

    def test_combine_slide_text_with_content(self):
        """Test combine_slide_text with a slide containing text shapes."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        shapes = slide.shapes
        txBox = shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
        tf = txBox.text_frame
        tf.text = "Hello"
        p = tf.add_paragraph()
        p.text = "world"

        combined_text = combine_slide_text(slide)
        self.assertEqual(combined_text, "Hello world")

    def test_combine_slide_text_no_content(self):
        """Test combine_slide_text with a slide containing no text shapes."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        combined_text = combine_slide_text(slide)
        self.assertEqual(combined_text, "")

    @patch('explainer.generate_prompt', return_value="test_prompt")
    @patch('explainer.fetch_explanation', new_callable=AsyncMock, return_value="test_explanation")
    async def test_process_slide_with_content(self, mock_fetch_explanation, mock_generate_prompt):
        """Test process_slide with a slide containing text."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        shapes = slide.shapes
        txBox = shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
        tf = txBox.text_frame
        tf.text = "Test content"

        client = AsyncMock()
        explanation = await process_slide(slide, client)
        self.assertEqual(explanation, "test_explanation")

    async def test_process_slide_no_content(self):
        """Test process_slide with a slide containing no text."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        client = AsyncMock()
        explanation = await process_slide(slide, client)
        self.assertEqual(explanation, "No text content")

    @patch('explainer.os.listdir', return_value=['test.pptx'])
    @patch('explainer.Presentation')
    @patch('explainer.process_slide', new_callable=AsyncMock, return_value="test_explanation")
    @patch('explainer.openai.AsyncClient')
    async def test_process_presentations(self, mock_openai_client, mock_process_slide, mock_presentation, mock_listdir):
        """Test process_presentations to ensure it processes files correctly."""
        mock_presentation_instance = mock_presentation.return_value
        mock_presentation_instance.slides = [AsyncMock()]

        await process_presentations()
        mock_openai_client.assert_called_once()
        mock_process_slide.assert_awaited()

if __name__ == '__main__':
    unittest.main()
