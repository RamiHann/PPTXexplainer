from pptx import Presentation


def extract_text_from_presentation(presentation_path):
    """
    Extract text from all slides in a PowerPoint presentation.

    Args:
    - presentation_path (str): Path to the PowerPoint presentation file.

    Returns:
    - List of strings, each representing the text from a slide.
    """
    presentation = Presentation(presentation_path)
    slides_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape_number, shape in enumerate(slide.shapes, start=1):
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Add slide and shape number for clarity
                    if shape_number == 1:
                        slide_text.append(f"Slide {slide_number}: ")
                    slide_text.append(text)

        if slide_text:
            # Combine all text parts into a single string
            combined_text = " ".join(slide_text).strip()
            slides_text.append(combined_text)

    return slides_text
