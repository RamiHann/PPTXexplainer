import os
import json
import asyncio
from venv import logger

from pptx import Presentation
import logging
from logging.handlers import TimedRotatingFileHandler
from dotenv import load_dotenv
import openai
from gpt_explainer import generate_prompt, fetch_explanation
from database import session, Upload
from datetime import datetime, timezone

# Constants
UPLOADS_FOLDER = 'uploads'
OUTPUTS_FOLDER = 'outputs'
LOGS_FOLDER = 'logs'
PROCESSING_LOG_FILE = os.path.join(LOGS_FOLDER, 'presentation_processing.log')
SUPPORTED_EXTENSIONS = {'.pptx'}

def configure_logging():
    """Configure logging for the application."""
    os.makedirs(LOGS_FOLDER, exist_ok=True)

    # Configure logging for slide processing
    logger = logging.getLogger()
    logger.setLevel(logging.DEBUG)

    # File handler
    file_handler = TimedRotatingFileHandler(PROCESSING_LOG_FILE, when='midnight', interval=1, backupCount=5)
    file_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    file_handler.setLevel(logging.INFO)

    # Stream handler (for terminal output)
    stream_handler = logging.StreamHandler()
    stream_handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    stream_handler.setLevel(logging.DEBUG)

    # Add handlers to the logger
    logger.addHandler(file_handler)
    logger.addHandler(stream_handler)

def load_env_variables():
    """Load environment variables from .env file."""
    load_dotenv()
    return os.getenv('OPENAI_API_KEY')

def combine_slide_text(slide):
    """Combine text from all shapes in a slide."""
    slide_texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
    return " ".join(slide_texts).strip()

async def process_slide(slide, client):
    """Process a single slide and return its explanation."""
    slide_text = combine_slide_text(slide)
    if slide_text:
        try:
            prompt = generate_prompt(slide_text)
            explanation = await fetch_explanation(client, prompt)
            return explanation
        except Exception as e:
            logger.error(f"Failed to process slide: {e}")
            return f"Failed to process slide: {e}"
    return "No text content"

async def process_presentations():
    """Process all new presentations in the UPLOADS_FOLDER."""
    openai_api_key = load_env_variables()
    client = openai.AsyncClient(api_key=openai_api_key)
    logger.info("Slide processing script started.")

<<<<<<< HEAD
    while True:
        pending_uploads = session.query(Upload).filter_by(status='pending').all()
        for upload in pending_uploads:
            pptx_path = os.path.join(UPLOADS_FOLDER, upload.filename)
            output_file = os.path.join(OUTPUTS_FOLDER, f"{upload.uid}.json")
=======
    while not asyncio.get_event_loop().is_closed():
        for filename in os.listdir(UPLOADS_FOLDER):
            if filename.endswith('.pptx'):
                pptx_path = os.path.join(UPLOADS_FOLDER, filename)
                output_file = os.path.join(OUTPUTS_FOLDER, f"{os.path.splitext(filename)[0]}.json")
>>>>>>> main

            if os.path.exists(output_file):
                continue  # Skip if already processed

<<<<<<< HEAD
            logger.info(f"Processing {upload.filename}...")
            prs = Presentation(pptx_path)
            tasks = []

            client = openai.AsyncClient(api_key=OPENAI_API_KEY)

            for slide in prs.slides:
                task = process_slide(slide, client)
                tasks.append(task)
=======
                logger.info(f"Processing {filename}...")
                presentation = Presentation(pptx_path)
                tasks = [process_slide(slide, client) for slide in presentation.slides]
>>>>>>> main

            explanations = await asyncio.gather(*tasks)

            with open(output_file, 'w') as f:
                json.dump(explanations, f, indent=4)

            # Update upload status in the database
            upload.status = 'done'
            upload.finish_time = datetime.utcnow()
            session.commit()

            logger.info(f"Processing {upload.filename} completed successfully.")

        await asyncio.sleep(10)  # Check for new uploads every 10 seconds

if __name__ == '__main__':
    configure_logging()
    try:
        asyncio.run(process_presentations())
    except KeyboardInterrupt:
        logger.info("Slide processing script ended due to keyboard interrupt.")
    except Exception as e:
        logger.error(f"Slide processing script ended with error: {e}")
